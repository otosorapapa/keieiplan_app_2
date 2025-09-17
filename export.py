"""Export utilities for generating deliverables from plan data."""
from __future__ import annotations

from io import BytesIO
from typing import Dict

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfgen import canvas

from utils import flatten_plan


def export_plan_to_excel(
    plan_data: Dict[str, Dict],
    income_statement: pd.DataFrame,
    forecast_df: pd.DataFrame,
    milestones_df: pd.DataFrame,
) -> bytes:
    """Generate an Excel workbook containing the plan artefacts."""

    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        flatten_plan(plan_data).to_excel(writer, sheet_name="Plan", index=False)
        income_statement.to_excel(writer, sheet_name="P&L", index=False)
        forecast_df.to_excel(writer, sheet_name="Forecast", index=False)
        milestones_df.to_excel(writer, sheet_name="Milestones", index=False)
    buffer.seek(0)
    return buffer.getvalue()


def export_plan_to_pdf(
    plan_data: Dict[str, Dict],
    narrative: str,
    income_statement: pd.DataFrame,
) -> bytes:
    """Create a lightweight PDF summary using ReportLab."""

    buffer = BytesIO()
    pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))
    c = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4
    margin = 40

    c.setFont("HeiseiKakuGo-W5", 16)
    c.drawString(margin, height - 60, f"{plan_data['overview']['company_name']} 経営計画要約")

    c.setFont("HeiseiKakuGo-W5", 12)
    y = height - 110
    for line in narrative.split("。"):
        if not line:
            continue
        c.drawString(margin, y, line + "。")
        y -= 18
        if y < margin:
            c.showPage()
            c.setFont("HeiseiKakuGo-W5", 12)
            y = height - margin

    c.setFont("HeiseiKakuGo-W5", 12)
    c.drawString(margin, y - 20, "主要KPI")
    y -= 50
    for _, row in income_statement.iterrows():
        text = f"{row['区分']}: {row['金額']:,}"
        c.drawString(margin, y, text)
        y -= 18
        if y < margin:
            c.showPage()
            c.setFont("HeiseiKakuGo-W5", 12)
            y = height - margin

    c.showPage()
    c.save()
    buffer.seek(0)
    return buffer.getvalue()


def export_plan_to_ppt(
    plan_data: Dict[str, Dict],
    narrative: str,
    income_statement: pd.DataFrame,
) -> bytes:
    """Produce a simple PowerPoint deck summarising the plan."""

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = plan_data["overview"].get("company_name", "経営計画")
    subtitle = slide.placeholders[1]
    subtitle.text = "統合事業計画ダッシュボード"

    # Strategy slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "戦略サマリー"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True
    for key in ["vision", "mission", "value_proposition", "target_market"]:
        paragraph = body.add_paragraph()
        paragraph.text = f"{key}: {plan_data['overview'].get(key, '')}"
        paragraph.font.size = Pt(18)

    # Financial slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "財務ハイライト"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    body.word_wrap = True
    for _, row in income_statement.head(5).iterrows():
        paragraph = body.add_paragraph()
        paragraph.text = f"{row['区分']}: {row['金額']:,}"
        paragraph.font.size = Pt(18)

    # Narrative slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "将来像"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for sentence in narrative.split("。"):
        if sentence.strip():
            paragraph = body.add_paragraph()
            paragraph.text = sentence.strip() + "。"
            paragraph.font.size = Pt(18)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
